#!/usr/bin/env python3
"""
Write alt text into each lecture deck's images and save a separate accessible
copy (<name>_accessible.pptx). Alt text comes from per-deck alt.json files
authored by viewing the extracted images.

Usage:
    python3 build_accessible_pptx.py <ALT_WORK_DIR>

<ALT_WORK_DIR> contains one subfolder per deck key (d1m, d1a, ...), each with
an alt.json mapping "slide<NNN>_pic<K>.<ext>" -> alt text ("" = decorative).

The accessible .pptx sets the PowerPoint "Description" (alt text) on every
picture, so the decks themselves are screen-reader friendly.
"""
import sys, os, json, re
from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

DECKS = {
 "d1m": "Day1_Probability/FW536_Day1_Morning_intro to probability.pptx",
 "d1a": "Day1_Probability/FW536_Day1_Afternoon Distributions I.pptx",
 "d2m": "Day2_GLM/FW536_Day2_Morning_ContinuousDistributions.pptx",
 "d2a": "Day2_GLM/FW536_Day2_Afternoon_GeneralizedLinear Models.pptx",
 "d3m": "Day3_MixedModels_ModelSelection/FW536_Day3_Morning_mixed effects.pptx",
 "d3a": "Day3_MixedModels_ModelSelection/FW536_Day3_Afternoon_hypothesis test and model selection.pptx",
 "d4m": "Day4_Likelihood_BayesI/FW536_Day4_Morning_maximum likelihod.pptx",
 "d4a": "Day4_Likelihood_BayesI/FW536_Day4_Afternoon_BayesI.pptx",
 "d5":  "Day5_BayesII/FW536_Day5_BayesII.pptx",
}


def set_descr(shape, text):
    cNvPr = shape._element.find('.//' + qn('p:cNvPr'))
    if cNvPr is None:
        return False
    cNvPr.set('descr', text or "")
    # a title helps some AT; keep short
    if text:
        cNvPr.set('title', text[:80])
    return True


def process(work):
    grand = {"decks": 0, "described": 0, "decorative": 0, "missing": 0}
    for key, rel in DECKS.items():
        altpath = os.path.join(work, key, "alt.json")
        if not os.path.exists(altpath):
            print(f"  [skip] {key}: no alt.json yet")
            continue
        alt = json.load(open(altpath))
        # normalize keys to slideNNN_picK (drop extension) for matching
        alt_by_sk = {}
        for fn, txt in alt.items():
            m = re.match(r"slide0*(\d+)_pic(\d+)", fn)
            if m:
                alt_by_sk[(int(m.group(1)), int(m.group(2)))] = txt
        src = os.path.join(ROOT, rel)
        prs = Presentation(src)
        described = decorative = missing = 0
        for si, slide in enumerate(prs.slides, 1):
            k = 0
            for sh in slide.shapes:
                if sh.shape_type == 13:  # picture
                    k += 1
                    txt = alt_by_sk.get((si, k))
                    if txt is None:
                        missing += 1
                        continue
                    set_descr(sh, txt)
                    if txt.strip():
                        described += 1
                    else:
                        decorative += 1
        out = re.sub(r"\.pptx$", "_accessible.pptx", src)
        prs.save(out)
        grand["decks"] += 1
        grand["described"] += described
        grand["decorative"] += decorative
        grand["missing"] += missing
        print(f"  {key}: described {described}, decorative {decorative}, missing {missing} -> {os.path.basename(out)}")
    print(f"\nTotal: {grand['decks']} decks, {grand['described']} described, "
          f"{grand['decorative']} decorative, {grand['missing']} still missing")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("usage: build_accessible_pptx.py <ALT_WORK_DIR>"); sys.exit(1)
    process(sys.argv[1])
