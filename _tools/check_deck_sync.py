#!/usr/bin/env python3
"""Check that each no-solution deck still agrees with the master it came from.

A no-solution deck is generated from its master by dropping whole slides and by
stripping answer content off slides that stay. So the invariant is one-directional:

    every paragraph on a no-solution slide must also appear on some master slide.

Text that exists only in the no-solution deck means the two have drifted -- almost
always because the master was edited after the no-solution deck was built, or
because an edit was applied to one file and not the other. That is the failure this
catches; it is invisible to a slide-count comparison, since counts stay equal.

Run from the repository root:

    python3 _tools/check_deck_sync.py

Exits 1 if any deck has drifted, so it can be used as a pre-teaching gate.
"""
import glob
import os
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE as T


def walk(shapes):
    """Yield every shape, descending into groups."""
    for sh in shapes:
        if sh.shape_type == T.GROUP:
            yield from walk(sh.shapes)
        else:
            yield sh


def paragraphs(slide):
    """Normalised non-empty paragraph strings on a slide."""
    out = []
    for sh in walk(slide.shapes):
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            text = ''.join(r.text for r in p.runs).strip()
            if text:
                out.append(' '.join(text.split()))
    return out


def pairs():
    for master in sorted(glob.glob('Day*/*_accessible.pptx')):
        if '_nosoln' in master:
            continue
        nosoln = master.replace('_accessible.pptx', '_nosoln_accessible.pptx')
        if os.path.exists(nosoln):
            yield master, nosoln


def main():
    drifted = 0
    for master, nosoln in pairs():
        master_slides = [set(paragraphs(s)) for s in Presentation(master).slides]
        nosoln_slides = list(Presentation(nosoln).slides)

        bad = []
        for i, slide in enumerate(nosoln_slides, 1):
            own = set(paragraphs(slide))
            if not any(own <= m for m in master_slides):
                orphans = [p for p in paragraphs(slide)
                           if not any(p in m for m in master_slides)]
                bad.append((i, orphans))

        name = os.path.basename(master).replace('FW536_', '').replace('_accessible.pptx', '')
        counts = f"{len(master_slides)} → {len(nosoln_slides)}"
        if bad:
            drifted += 1
            print(f"DRIFTED  {name}  ({counts})")
            for i, orphans in bad:
                print(f"           no-solution slide {i} carries text the master does not:")
                for line in orphans[:4]:
                    print(f"             {line[:100]}")
        else:
            print(f"ok       {name}  ({counts})")

    print()
    if drifted:
        print(f"{drifted} deck(s) drifted. Re-apply the edit to both files, "
              f"or rebuild the no-solution deck from its master.")
        return 1
    print("All no-solution decks agree with their masters.")
    return 0


if __name__ == '__main__':
    sys.exit(main())
