#!/usr/bin/env python3
"""
Generate accessible HTML lecture notes from the lecture decks. If an
"<name>_accessible.pptx" (with alt text in each picture's Description field)
exists, it is used and the alt text becomes each figure's description; images
with empty alt are treated as decorative; images with no alt are flagged
"needs description". Writes into accessible/Day*/ and refreshes
FIGURE_DESCRIPTIONS_TODO.md.

Run after build_accessible_pptx.py.
"""
import os, re, html as H
from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ACC = os.path.join(ROOT, "accessible")

DECKS = [
 ("Day1_Probability", "FW536_Day1_Morning_intro to probability.pptx", "lecture_day1_morning.html", "Day 1 Morning", "Introduction to probability"),
 ("Day1_Probability", "FW536_Day1_Afternoon Distributions I.pptx", "lecture_day1_afternoon.html", "Day 1 Afternoon", "Discrete distributions"),
 ("Day2_GLM", "FW536_Day2_Morning_ContinuousDistributions.pptx", "lecture_day2_morning.html", "Day 2 Morning", "Continuous distributions"),
 ("Day2_GLM", "FW536_Day2_Afternoon_GeneralizedLinear Models.pptx", "lecture_day2_afternoon.html", "Day 2 Afternoon", "Generalized linear models"),
 ("Day3_MixedModels_ModelSelection", "FW536_Day3_Morning_mixed effects.pptx", "lecture_day3_morning.html", "Day 3 Morning", "Mixed effects models"),
 ("Day3_MixedModels_ModelSelection", "FW536_Day3_Afternoon_hypothesis test and model selection.pptx", "lecture_day3_afternoon.html", "Day 3 Afternoon", "Hypothesis tests and model selection"),
 ("Day4_Likelihood_BayesI", "FW536_Day4_Morning_maximum likelihod.pptx", "lecture_day4_morning.html", "Day 4 Morning", "Maximum likelihood"),
 ("Day4_Likelihood_BayesI", "FW536_Day4_Afternoon_BayesI.pptx", "lecture_day4_afternoon.html", "Day 4 Afternoon", "Bayesian statistics I"),
 ("Day5_BayesII", "FW536_Day5_BayesII.pptx", "lecture_day5.html", "Day 5", "Bayesian statistics II"),
]

HEAD = """<!DOCTYPE html>
<html lang="en" data-a11y="done">
<head>
<meta charset="UTF-8" />
<title>{title} lecture notes | FW 536</title>
<meta name="viewport" content="width=device-width, initial-scale=1" />
<link rel="stylesheet" href="../_Shared/css/accessible.css" />
<script>
MathJax = {{ tex: {{ inlineMath: [['$','$'], ['\\\\(','\\\\)']] }},
  options: {{ menuOptions: {{ settings: {{ assistiveMml: true, explorer: true }} }} }},
  loader: {{ load: ['a11y/assistive-mml', 'a11y/explorer'] }} }};
</script>
<script src="https://cdn.jsdelivr.net/npm/mathjax@3/es5/tex-chtml.js" id="MathJax-script"></script>
</head>
<body>
<a class="skip-link" href="#main-content">Skip to main content</a>
<div class="a11y-bar" role="region" aria-label="Accessibility">Accessible version &middot; <a href="../accessibility.html">How to use this site with a screen reader</a> &middot; <a href="../index.html">Course home</a></div>
"""
FOOT = """
<footer>FW 536 &middot; Statistical Modeling for Ecology and Conservation &middot; Oregon State University &middot; Accessible lecture notes</footer>
</body></html>"""


def esc(s):
    return H.escape(s).replace("\n", " ").strip()

def alt_of(shape):
    cNvPr = shape._element.find('.//' + qn('p:cNvPr'))
    if cNvPr is None:
        return None
    d = cNvPr.get('descr')
    return d  # None if never set; "" if explicitly decorative

def deck_path(folder, pptx):
    acc = os.path.join(ROOT, folder, re.sub(r"\.pptx$", "_accessible.pptx", pptx))
    return acc if os.path.exists(acc) else os.path.join(ROOT, folder, pptx)

def extract(slide):
    title = None
    if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
        t = slide.shapes.title.text_frame.text.strip()
        if t: title = t
    blocks = []; pics = 0
    for shape in slide.shapes:
        if shape == slide.shapes.title: continue
        if shape.shape_type == 13:
            pics += 1; blocks.append({"type": "figure", "n": pics, "alt": alt_of(shape)})
        elif shape.has_table:
            blocks.append({"type": "table", "rows": [[c.text.strip() for c in r.cells] for r in shape.table.rows]})
        elif shape.has_text_frame:
            paras = []
            for p in shape.text_frame.paragraphs:
                txt = "".join(r.text for r in p.runs).strip() or p.text.strip()
                if txt: paras.append((p.level or 0, txt))
            if paras: blocks.append({"type": "text", "paras": paras})
    return title, blocks

def render_blocks(blocks, slide_no, missing_log, label):
    out = []
    for b in blocks:
        if b["type"] == "text":
            if len(b["paras"]) == 1 and b["paras"][0][0] == 0:
                out.append(f"<p>{esc(b['paras'][0][1])}</p>")
            else:
                out.append("<ul>")
                for lvl, txt in b["paras"]:
                    pad = f' style="margin-left:{lvl*1.2}em"' if lvl else ""
                    out.append(f"<li{pad}>{esc(txt)}</li>")
                out.append("</ul>")
        elif b["type"] == "table":
            rows = b["rows"]
            if not rows: continue
            out.append("<table><thead><tr>" + "".join(f'<th scope="col">{esc(c)}</th>' for c in rows[0]) + "</tr></thead>")
            if len(rows) > 1:
                out.append("<tbody>")
                for r in rows[1:]:
                    out.append("<tr>" + "".join(f"<td>{esc(c)}</td>" for c in r) + "</tr>")
                out.append("</tbody>")
            out.append("</table>")
        elif b["type"] == "figure":
            k = b["n"]; alt = b["alt"]
            if alt is None:
                missing_log.append((label, slide_no, k))
                out.append(f'<p class="callout warn" role="note"><span class="callout-label">Figure {slide_no}.{k} — needs description</span> '
                           f'This slide contains an image not yet described in text. <em>Instructor: add a description.</em></p>')
            elif alt.strip() == "":
                out.append(f'<p class="sr-only">Figure {slide_no}.{k}: decorative image, no description needed.</p>')
            else:
                out.append(f'<figure role="group"><figcaption><strong>Figure {slide_no}.{k}.</strong> {esc(alt)}</figcaption></figure>')
    return "\n".join(out)

def build(folder, pptx, out_name, label, subtitle, missing_log):
    prs = Presentation(deck_path(folder, pptx))
    parts = [HEAD.format(title=f"{label} — {subtitle}")]
    parts.append('<aside class="sidebar">')
    parts.append(f'<p class="brand"><a href="../index.html">FW 536 · {label}</a></p>')
    parts.append(f'<div class="subtitle">Lecture notes: {esc(subtitle)}</div>')
    parts.append('<nav aria-label="Slides in this lecture"><div class="section">Slides</div><ul>')
    body = ['<main class="content" id="main-content" tabindex="-1">',
            '<header class="lab-header">',
            f'<div class="eyebrow">{esc(label)} · Accessible lecture notes</div>',
            f'<h1>{esc(subtitle)}</h1>',
            '<p class="lede">A screen-reader-friendly transcript of the lecture slides: all slide text, tables, figure descriptions, and speaker notes.</p>',
            '</header>',
            '<div class="callout key"><p class="callout-label">About these notes</p>'
            'Each section is one slide. Figures carry text descriptions where available. '
            'For the concepts with worked examples, also read the '
            '<a href="plain_language_summary.html">plain-language summary</a> and '
            '<a href="lab.html">lab</a> for this session.</div>']
    for i, slide in enumerate(prs.slides, 1):
        title, blocks = extract(slide)
        anchor = f"slide-{i}"
        nav_label = title if title else f"Slide {i}"
        parts.append(f'<li><a href="#{anchor}">{i}. {esc(nav_label)}</a></li>')
        body.append(f'<section aria-labelledby="{anchor}-h">')
        head_txt = f"Slide {i}" + (f": {esc(title)}" if title else "")
        body.append(f'<h2 id="{anchor}-h"><span id="{anchor}"></span>{head_txt}</h2>')
        r = render_blocks(blocks, i, missing_log, label)
        body.append(r if r.strip() else "<p><em>(No text on this slide.)</em></p>")
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                body.append('<details class="answer"><summary>Speaker notes for this slide</summary>'
                            f'<div class="answer-body"><p>{esc(notes)}</p></div></details>')
        body.append('</section>')
    parts.append('</ul></nav></aside>')
    parts.extend(body); parts.append('</main>'); parts.append(FOOT)
    with open(os.path.join(ACC, folder, out_name), "w", encoding="utf-8") as f:
        f.write("\n".join(parts))
    return len(prs.slides)

def main():
    missing = []
    for folder, pptx, out_name, label, subtitle in DECKS:
        build(folder, pptx, out_name, label, subtitle, missing)
        print(f"  wrote accessible/{folder}/{out_name}")
    todo = ["# Figure descriptions still needed (accessible lecture notes)\n"]
    if not missing:
        todo.append("All lecture figures now have alt-text descriptions (or are marked decorative). Nothing outstanding.\n")
    else:
        todo.append(f"{len(missing)} figures still need a description:\n")
        cur = None
        for label, s, k in missing:
            if label != cur: todo.append(f"\n### {label}"); cur = label
            todo.append(f"- [ ] Figure {s}.{k} (slide {s})")
    open(os.path.join(ACC, "FIGURE_DESCRIPTIONS_TODO.md"), "w").write("\n".join(todo) + "\n")
    print(f"\n{len(missing)} figures still missing descriptions")

if __name__ == "__main__":
    main()
